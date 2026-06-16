"""
AeroVision AI - Tata Technologies InnoVent Presentation Generator
Generates a professional 19-slide PPTX with white & dark blue theme.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData
import os

# ─── Theme Colors ─────────────────────────────────────────────────────────────
DARK_BLUE    = RGBColor(0x00, 0x1F, 0x3F)   # #001F3F
NAVY_BLUE    = RGBColor(0x00, 0x2B, 0x5C)   # #002B5C
MEDIUM_BLUE  = RGBColor(0x00, 0x74, 0xD9)   # #0074D9
LIGHT_BLUE   = RGBColor(0x7F, 0xDB, 0xFF)   # #7FDBFF
SKY_BLUE     = RGBColor(0x39, 0xCC, 0xCC)   # #39CCCC
ACCENT_BLUE  = RGBColor(0x00, 0x9E, 0xEB)   # #009EEB
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE    = RGBColor(0xF0, 0xF4, 0xF8)   # #F0F4F8
LIGHT_GRAY   = RGBColor(0xE8, 0xEC, 0xF0)   # #E8ECF0
DARK_GRAY    = RGBColor(0x33, 0x33, 0x33)   # #333333
ORANGE       = RGBColor(0xFF, 0x85, 0x1B)   # #FF851B
GREEN        = RGBColor(0x2E, 0xCC, 0x71)   # #2ECC71
RED_ACCENT   = RGBColor(0xE7, 0x4C, 0x3C)   # #E74C3C

# Slide dimensions (16:9)
SLIDE_WIDTH  = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT

# ─── Helper Functions ─────────────────────────────────────────────────────────

def add_dark_bg(slide):
    """Add dark blue background to slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BLUE

def add_white_bg(slide):
    """Add white background to slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

def add_gradient_bg(slide):
    """Add a solid dark navy background."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = NAVY_BLUE

def add_shape(slide, shape_type, left, top, width, height, fill_color=None, line_color=None, line_width=None):
    """Add a shape with optional fill and line."""
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape

def add_textbox(slide, left, top, width, height, text, font_size=18, font_color=WHITE, 
                bold=False, alignment=PP_ALIGN.LEFT, font_name='Calibri', line_spacing=1.2):
    """Add a text box with specified formatting."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(0)
    p.line_spacing = Pt(font_size * line_spacing)
    return txBox

def add_multi_text(slide, left, top, width, height, lines, font_size=16, font_color=WHITE, 
                   font_name='Calibri', bold_first=False, bullet=False, line_spacing=1.5):
    """Add a text box with multiple paragraphs."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        prefix = "●  " if bullet else ""
        p.text = prefix + line
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.name = font_name
        p.font.bold = (bold_first and i == 0)
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(4)
        p.line_spacing = Pt(font_size * line_spacing)
    return txBox

def add_accent_bar(slide, left, top, width=Inches(0.08), height=Inches(0.6), color=MEDIUM_BLUE):
    """Add a vertical accent bar."""
    return add_shape(slide, MSO_SHAPE.RECTANGLE, left, top, width, height, fill_color=color)

def add_horizontal_line(slide, left, top, width, color=MEDIUM_BLUE, thickness=2):
    """Add a horizontal line."""
    return add_shape(slide, MSO_SHAPE.RECTANGLE, left, top, width, Pt(thickness), fill_color=color)

def add_card(slide, left, top, width, height, fill_color=NAVY_BLUE, border_color=None):
    """Add a card/panel shape."""
    shape = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height, 
                      fill_color=fill_color, line_color=border_color, line_width=1.5 if border_color else None)
    return shape

def add_slide_number(slide, num, total=19):
    """Add slide number at bottom right."""
    add_textbox(slide, Inches(11.5), Inches(7.0), Inches(1.5), Inches(0.4),
                f"{num} / {total}", font_size=10, font_color=RGBColor(0x66, 0x77, 0x88),
                alignment=PP_ALIGN.RIGHT)

def add_top_bar(slide):
    """Add decorative top accent bar."""
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), fill_color=MEDIUM_BLUE)

def add_bottom_bar(slide):
    """Add decorative bottom bar."""
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.35), SLIDE_WIDTH, Inches(0.15), fill_color=MEDIUM_BLUE)

def add_section_header(slide, title, subtitle=None, dark=True):
    """Add a section title with optional subtitle."""
    fc = WHITE if dark else DARK_BLUE
    sc = LIGHT_BLUE if dark else MEDIUM_BLUE
    add_accent_bar(slide, Inches(0.8), Inches(0.5), Inches(0.08), Inches(0.55), color=MEDIUM_BLUE)
    add_textbox(slide, Inches(1.05), Inches(0.4), Inches(10), Inches(0.7), title, 
                font_size=32, font_color=fc, bold=True, font_name='Calibri')
    if subtitle:
        add_textbox(slide, Inches(1.05), Inches(1.0), Inches(10), Inches(0.4), subtitle,
                    font_size=14, font_color=sc, font_name='Calibri')
    add_horizontal_line(slide, Inches(0.8), Inches(1.35), Inches(11.5), color=MEDIUM_BLUE, thickness=2)

def add_white_section_header(slide, title, subtitle=None):
    """Section header for white background slides."""
    add_section_header(slide, title, subtitle, dark=False)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1: TITLE SLIDE
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
add_dark_bg(slide)

# Top decorative element
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.08), fill_color=MEDIUM_BLUE)

# Left vertical blue accent
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.15), SLIDE_HEIGHT, fill_color=MEDIUM_BLUE)

# Main title area
add_textbox(slide, Inches(1.5), Inches(1.2), Inches(10), Inches(1.0),
            "AEROVISION AI", font_size=52, font_color=WHITE, bold=True, font_name='Calibri')

add_textbox(slide, Inches(1.5), Inches(2.2), Inches(10), Inches(0.6),
            "AI-Powered Aircraft Surface Defect Detection", font_size=24, 
            font_color=LIGHT_BLUE, font_name='Calibri')

# Separator line
add_horizontal_line(slide, Inches(1.5), Inches(3.0), Inches(4), color=MEDIUM_BLUE, thickness=3)

# Team & Details
details = [
    ("Team", "Team Nexus Innovation"),
    ("Team Lead", "Mugesh Kumar M"),
    ("Members", "Karthikeyan S  |  Gokulakrishnan S  |  Pradeep Kumar R"),
    ("College", "K.S.Rangasamy College of Technology, Tiruchengode"),
    ("Industry", "Aerospace & Aviation"),
    ("Category", "AI-at-the-Edge / Real-Time Detection"),
]

y_pos = 3.4
for label, value in details:
    add_textbox(slide, Inches(1.5), Inches(y_pos), Inches(2.5), Inches(0.35),
                label.upper(), font_size=11, font_color=MEDIUM_BLUE, bold=True, font_name='Calibri')
    add_textbox(slide, Inches(3.8), Inches(y_pos), Inches(8), Inches(0.35),
                value, font_size=14, font_color=WHITE, font_name='Calibri')
    y_pos += 0.42

# Competition badge
badge = add_card(slide, Inches(9), Inches(1.2), Inches(3.5), Inches(1.5), fill_color=NAVY_BLUE, border_color=MEDIUM_BLUE)
add_textbox(slide, Inches(9.2), Inches(1.35), Inches(3.1), Inches(0.4),
            "TATA TECHNOLOGIES", font_size=14, font_color=MEDIUM_BLUE, bold=True, 
            font_name='Calibri', alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(9.2), Inches(1.75), Inches(3.1), Inches(0.5),
            "InnoVent 2026", font_size=28, font_color=WHITE, bold=True, 
            font_name='Calibri', alignment=PP_ALIGN.CENTER)

# Bottom bar
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.35), SLIDE_WIDTH, Inches(0.15), fill_color=MEDIUM_BLUE)
add_slide_number(slide, 1)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2: INTRODUCTION
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_top_bar(slide)
add_section_header(slide, "TEAM INTRODUCTION", "Meet the Team Behind AeroVision AI")

# Team member cards
members = [
    ("Mugesh Kumar M", "Team Lead & AI/ML Developer", 
     "Leads project architecture, designs and trains YOLOv8 models, manages ML pipeline and edge deployment strategy."),
    ("Karthikeyan S", "Edge Computing & IoT Specialist", 
     "Responsible for edge hardware setup, NVIDIA Jetson/Raspberry Pi deployment, and IoT sensor integration."),
    ("Gokulakrishnan S", "Computer Vision Engineer", 
     "Develops computer vision algorithms, handles image preprocessing, and implements OpenCV pipelines."),
    ("Pradeep Kumar R", "Full-Stack & Data Engineer", 
     "Builds Flask web dashboard, manages database, develops REST APIs, and handles data pipelines."),
]

x_positions = [0.5, 3.6, 6.7, 9.8]
for i, (name, role, desc) in enumerate(members):
    x = x_positions[i]
    # Card background
    add_card(slide, Inches(x), Inches(1.8), Inches(2.8), Inches(4.2), fill_color=NAVY_BLUE, border_color=MEDIUM_BLUE)
    # Icon circle
    add_shape(slide, MSO_SHAPE.OVAL, Inches(x + 0.9), Inches(2.0), Inches(1.0), Inches(1.0), fill_color=MEDIUM_BLUE)
    add_textbox(slide, Inches(x + 0.9), Inches(2.15), Inches(1.0), Inches(0.7),
                "👤", font_size=32, font_color=WHITE, alignment=PP_ALIGN.CENTER)
    # Name
    add_textbox(slide, Inches(x + 0.1), Inches(3.15), Inches(2.6), Inches(0.4),
                name, font_size=15, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER, font_name='Calibri')
    # Role
    add_textbox(slide, Inches(x + 0.1), Inches(3.55), Inches(2.6), Inches(0.45),
                role, font_size=10, font_color=LIGHT_BLUE, alignment=PP_ALIGN.CENTER, font_name='Calibri')
    # Description
    add_textbox(slide, Inches(x + 0.2), Inches(4.1), Inches(2.4), Inches(1.6),
                desc, font_size=10, font_color=OFF_WHITE, alignment=PP_ALIGN.CENTER, font_name='Calibri')

# Project overview box
add_card(slide, Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.9), fill_color=NAVY_BLUE, border_color=MEDIUM_BLUE)
add_textbox(slide, Inches(0.7), Inches(6.3), Inches(12), Inches(0.7),
            "AeroVision AI is an AI-powered edge computing system for real-time aircraft surface defect detection, "
            "using a custom YOLOv8 model with ResCBAM attention and GhostConv for lightweight, high-accuracy inspection.",
            font_size=13, font_color=OFF_WHITE, alignment=PP_ALIGN.CENTER, font_name='Calibri')

add_bottom_bar(slide)
add_slide_number(slide, 2)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3: PROBLEM STATEMENT
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_white_bg(slide)
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), fill_color=DARK_BLUE)
add_white_section_header(slide, "PROBLEM STATEMENT", "Why Aircraft Surface Inspection Needs a Revolution")

# Problem cards
problems = [
    ("⏱️", "Slow Inspection", "4+ hours per aircraft using manual visual inspection methods"),
    ("👁️", "Low Accuracy", "Only 85% detection rate due to human fatigue and subjective judgment"),
    ("💰", "High Cost", "$200K+ annual facility cost with $100-300 per single inspection"),
    ("📉", "Low Throughput", "Only 6 aircraft inspected per day, creating bottlenecks"),
]

for i, (icon, title, desc) in enumerate(problems):
    x = Inches(0.5 + i * 3.15)
    add_card(slide, x, Inches(1.7), Inches(2.85), Inches(2.3), fill_color=DARK_BLUE)
    add_textbox(slide, x, Inches(1.85), Inches(2.85), Inches(0.5),
                icon, font_size=36, font_color=WHITE, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5 + i * 3.15 + 0.15), Inches(2.45), Inches(2.55), Inches(0.35),
                title, font_size=16, font_color=LIGHT_BLUE, bold=True, alignment=PP_ALIGN.CENTER, font_name='Calibri')
    add_textbox(slide, Inches(0.5 + i * 3.15 + 0.15), Inches(2.85), Inches(2.55), Inches(0.9),
                desc, font_size=12, font_color=OFF_WHITE, alignment=PP_ALIGN.CENTER, font_name='Calibri')

# Current limitations section
add_card(slide, Inches(0.5), Inches(4.3), Inches(12.3), Inches(2.8), fill_color=OFF_WHITE)
add_textbox(slide, Inches(0.8), Inches(4.4), Inches(5), Inches(0.4),
            "CURRENT LIMITATIONS", font_size=18, font_color=DARK_BLUE, bold=True, font_name='Calibri')

limitations_left = [
    "Manual inspection is inconsistent across inspectors",
    "Human fatigue leads to missed critical defects",
    "No real-time feedback during inspection process",
]
limitations_right = [
    "Existing AI solutions cost $200K-$250K/year",
    "Proprietary vendor lock-in with limited flexibility",
    "Slow inference speed (500ms-1s) not suitable for real-time",
]

add_multi_text(slide, Inches(0.8), Inches(4.9), Inches(5.5), Inches(2.0), limitations_left,
               font_size=13, font_color=DARK_GRAY, bullet=True, line_spacing=1.6)
add_multi_text(slide, Inches(6.5), Inches(4.9), Inches(6), Inches(2.0), limitations_right,
               font_size=13, font_color=DARK_GRAY, bullet=True, line_spacing=1.6)

add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.35), SLIDE_WIDTH, Inches(0.15), fill_color=DARK_BLUE)
add_slide_number(slide, 3)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4: MARKET SIZE AND FUTURE POTENTIAL
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_top_bar(slide)
add_section_header(slide, "MARKET SIZE & FUTURE POTENTIAL", "A Rapidly Growing Industry Opportunity")

# Market stats cards
stats = [
    ("$15B", "Aircraft Maintenance\nInspection Market", MEDIUM_BLUE),
    ("8.5%", "CAGR Growth\nto 2030", ACCENT_BLUE),
    ("$26B", "Projected Market\nby 2030", SKY_BLUE),
    ("$500M", "Addressable Market\n(First 5 Years)", GREEN),
]

for i, (value, label, color) in enumerate(stats):
    x = Inches(0.5 + i * 3.15)
    add_card(slide, x, Inches(1.7), Inches(2.85), Inches(1.8), fill_color=NAVY_BLUE, border_color=color)
    add_textbox(slide, x, Inches(1.85), Inches(2.85), Inches(0.65),
                value, font_size=36, font_color=color, bold=True, alignment=PP_ALIGN.CENTER, font_name='Calibri')
    add_textbox(slide, x, Inches(2.55), Inches(2.85), Inches(0.75),
                label, font_size=12, font_color=OFF_WHITE, alignment=PP_ALIGN.CENTER, font_name='Calibri')

# Market details
add_card(slide, Inches(0.5), Inches(3.8), Inches(6), Inches(3.2), fill_color=NAVY_BLUE, border_color=MEDIUM_BLUE)
add_textbox(slide, Inches(0.8), Inches(3.95), Inches(5.5), Inches(0.4),
            "MARKET SEGMENTS", font_size=16, font_color=LIGHT_BLUE, bold=True, font_name='Calibri')

segments = [
    "Global Aerospace NDT Market: $4.2B (2024) → $7.8B (2030)",
    "AI-based Quality Inspection: $1.2B → $4.5B by 2028",
    "Global Aerospace MRO Market: $90B+ annually",
    "Defense & Military Aviation Inspection: High growth",
    "Emerging Markets (India, Asia-Pacific): Fastest growth",
]
add_multi_text(slide, Inches(0.8), Inches(4.4), Inches(5.5), Inches(2.4), segments,
               font_size=13, font_color=OFF_WHITE, bullet=True, line_spacing=1.5)

# TAM/SAM/SOM
add_card(slide, Inches(6.8), Inches(3.8), Inches(5.9), Inches(3.2), fill_color=NAVY_BLUE, border_color=MEDIUM_BLUE)
add_textbox(slide, Inches(7.1), Inches(3.95), Inches(5.4), Inches(0.4),
            "BUSINESS OPPORTUNITY", font_size=16, font_color=LIGHT_BLUE, bold=True, font_name='Calibri')

tam_data = [
    "TAM (Total Addressable Market): $4.5B",
    "SAM (Serviceable Available Market): $1.2B",
    "SOM (Serviceable Obtainable Market): $50M",
    "Break-even Period: 6-12 months",
    "5-Year ROI for 500-Aircraft Fleet: $486K (548%)",
]
add_multi_text(slide, Inches(7.1), Inches(4.4), Inches(5.4), Inches(2.4), tam_data,
               font_size=13, font_color=OFF_WHITE, bullet=True, line_spacing=1.5)

add_bottom_bar(slide)
add_slide_number(slide, 4)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5: PROJECT OBJECTIVE
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_white_bg(slide)
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), fill_color=DARK_BLUE)
add_white_section_header(slide, "PROJECT OBJECTIVE", "Defining Clear Goals for Impact")

# Main objective
add_card(slide, Inches(0.5), Inches(1.7), Inches(12.3), Inches(1.5), fill_color=DARK_BLUE)
add_textbox(slide, Inches(0.8), Inches(1.8), Inches(1.5), Inches(0.4),
            "MAIN OBJECTIVE", font_size=11, font_color=MEDIUM_BLUE, bold=True, font_name='Calibri')
add_textbox(slide, Inches(0.8), Inches(2.2), Inches(11.5), Inches(0.7),
            "Develop an AI-powered edge computing system for real-time aircraft surface defect detection that "
            "achieves superior accuracy, speed, and cost-efficiency compared to manual inspection methods.",
            font_size=16, font_color=WHITE, font_name='Calibri')

# Expected outcomes (3 columns)
outcomes = [
    ("🎯", "DETECTION ACCURACY", "Achieve ≥94% mAP@50 defect\ndetection accuracy across all\n6 defect categories"),
    ("⚡", "REAL-TIME PROCESSING", "Process images in <100ms\nper frame enabling 30 FPS\nreal-time inspection"),
    ("💡", "COST REDUCTION", "Reduce inspection costs by\n60% and inspection time by\n70% vs manual methods"),
]

for i, (icon, title, desc) in enumerate(outcomes):
    x = Inches(0.5 + i * 4.15)
    add_card(slide, x, Inches(3.5), Inches(3.85), Inches(2.5), fill_color=OFF_WHITE)
    add_textbox(slide, x, Inches(3.6), Inches(3.85), Inches(0.55),
                icon, font_size=36, font_color=DARK_BLUE, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5 + i * 4.15 + 0.2), Inches(4.2), Inches(3.45), Inches(0.35),
                title, font_size=13, font_color=DARK_BLUE, bold=True, alignment=PP_ALIGN.CENTER, font_name='Calibri')
    add_textbox(slide, Inches(0.5 + i * 4.15 + 0.2), Inches(4.65), Inches(3.45), Inches(1.2),
                desc, font_size=12, font_color=DARK_GRAY, alignment=PP_ALIGN.CENTER, font_name='Calibri')

# Key goals
add_textbox(slide, Inches(0.8), Inches(6.3), Inches(3), Inches(0.3),
            "KEY GOALS", font_size=14, font_color=DARK_BLUE, bold=True, font_name='Calibri')
goals = ["Zero missed critical defects", "Edge-deployable lightweight model (<25MB)", 
         "Web-based monitoring dashboard", "Scalable multi-camera support"]
add_multi_text(slide, Inches(0.8), Inches(6.6), Inches(12), Inches(0.8), goals,
               font_size=12, font_color=DARK_GRAY, bullet=True, line_spacing=1.3)

add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.35), SLIDE_WIDTH, Inches(0.15), fill_color=DARK_BLUE)
add_slide_number(slide, 5)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6: PROPOSED SOLUTION
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_top_bar(slide)
add_section_header(slide, "PROPOSED SOLUTION", "AeroVision AI — How It Works")

# 3-Step Process
steps = [
    ("01", "IMAGE ACQUISITION", "High-resolution surface images\ncaptured via UAV, handheld\ncamera, or automated system\nat 640×640 resolution", "📷"),
    ("02", "AI PROCESSING", "Custom YOLOv8 neural network\nwith ResCBAM attention\nprocesses image in 11-79ms\non GPU-accelerated edge device", "🧠"),
    ("03", "INSTANT RESULTS", "Visual report with highlighted\ndefects, confidence scores,\nautomatic alerts, and\ndetailed classification", "📊"),
]

for i, (num, title, desc, icon) in enumerate(steps):
    x = Inches(0.5 + i * 4.15)
    add_card(slide, x, Inches(1.7), Inches(3.85), Inches(3.5), fill_color=NAVY_BLUE, border_color=MEDIUM_BLUE)
    # Step number circle
    add_shape(slide, MSO_SHAPE.OVAL, Inches(0.5 + i * 4.15 + 1.3), Inches(1.9), Inches(1.2), Inches(1.2), fill_color=MEDIUM_BLUE)
    add_textbox(slide, Inches(0.5 + i * 4.15 + 1.3), Inches(2.05), Inches(1.2), Inches(0.9),
                num, font_size=30, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER, font_name='Calibri')
    # Title
    add_textbox(slide, Inches(0.5 + i * 4.15 + 0.2), Inches(3.25), Inches(3.45), Inches(0.35),
                title, font_size=15, font_color=LIGHT_BLUE, bold=True, alignment=PP_ALIGN.CENTER, font_name='Calibri')
    # Description
    add_textbox(slide, Inches(0.5 + i * 4.15 + 0.2), Inches(3.7), Inches(3.45), Inches(1.3),
                desc, font_size=12, font_color=OFF_WHITE, alignment=PP_ALIGN.CENTER, font_name='Calibri')

# Arrows between steps
for i in range(2):
    x = Inches(4.15 + i * 4.15)
    add_textbox(slide, x, Inches(2.7), Inches(0.7), Inches(0.5),
                "➤", font_size=28, font_color=MEDIUM_BLUE, alignment=PP_ALIGN.CENTER)

# 6 Defect types
add_card(slide, Inches(0.5), Inches(5.5), Inches(12.3), Inches(1.6), fill_color=NAVY_BLUE, border_color=MEDIUM_BLUE)
add_textbox(slide, Inches(0.8), Inches(5.6), Inches(5), Inches(0.3),
            "6 DEFECT TYPES DETECTED", font_size=14, font_color=LIGHT_BLUE, bold=True, font_name='Calibri')

defects = [
    ("Crack", "🔴"), ("Inclusion", "🟣"), ("Corrosion", "🟠"),
    ("Surface Wear", "🟢"), ("Delamination", "🟤"), ("Scratch", "🟡"),
]

for i, (name, color) in enumerate(defects):
    x = Inches(0.8 + i * 2.05)
    add_textbox(slide, x, Inches(6.05), Inches(1.8), Inches(0.35),
                f"{color}  {name}", font_size=13, font_color=WHITE, font_name='Calibri')

add_bottom_bar(slide)
add_slide_number(slide, 6)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7: INNOVATION AND NOVELTY
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_white_bg(slide)
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), fill_color=DARK_BLUE)
add_white_section_header(slide, "INNOVATION & NOVELTY", "What Makes AeroVision AI Unique")

# 3 Custom architecture improvements
innovations = [
    ("ResCBAM Attention", 
     "Residual connection with Convolutional Block Attention Module (CBAM) for superior feature extraction and reduced overfitting. Focuses on relevant defect features.",
     "🔬"),
    ("GhostConv Replacement", 
     "Lightweight ghost convolutions in backbone achieve 40% parameter reduction without sacrificing detection performance. Enables edge deployment.",
     "⚡"),
    ("WIOU Loss Function", 
     "Improved Weighted IoU loss function handles imbalanced data and prevents small target defects from dominating the training process.",
     "🎯"),
]

for i, (title, desc, icon) in enumerate(innovations):
    x = Inches(0.5 + i * 4.15)
    add_card(slide, x, Inches(1.7), Inches(3.85), Inches(2.5), fill_color=DARK_BLUE)
    add_textbox(slide, x, Inches(1.85), Inches(3.85), Inches(0.55),
                icon, font_size=36, font_color=WHITE, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5 + i * 4.15 + 0.2), Inches(2.45), Inches(3.45), Inches(0.35),
                title, font_size=15, font_color=LIGHT_BLUE, bold=True, alignment=PP_ALIGN.CENTER, font_name='Calibri')
    add_textbox(slide, Inches(0.5 + i * 4.15 + 0.2), Inches(2.85), Inches(3.45), Inches(1.2),
                desc, font_size=11, font_color=OFF_WHITE, alignment=PP_ALIGN.CENTER, font_name='Calibri')

# Competitive advantages
add_card(slide, Inches(0.5), Inches(4.5), Inches(6), Inches(2.6), fill_color=OFF_WHITE)
add_textbox(slide, Inches(0.8), Inches(4.6), Inches(5.5), Inches(0.35),
            "COMPETITIVE ADVANTAGES", font_size=15, font_color=DARK_BLUE, bold=True, font_name='Calibri')
advantages = [
    "4.05M parameters only — 42% lighter than YOLOv5s",
    "79.2% mAP@50 — best among all compared models",
    "Real-time edge inference without cloud dependency",
    "Multi-defect classification in a single pass",
    "Open architecture — no vendor lock-in",
]
add_multi_text(slide, Inches(0.8), Inches(5.05), Inches(5.5), Inches(1.8), advantages,
               font_size=12, font_color=DARK_GRAY, bullet=True, line_spacing=1.5)

# AI at the Edge
add_card(slide, Inches(6.8), Inches(4.5), Inches(6), Inches(2.6), fill_color=OFF_WHITE)
add_textbox(slide, Inches(7.1), Inches(4.6), Inches(5.5), Inches(0.35),
            "AI-AT-THE-EDGE HIGHLIGHTS", font_size=15, font_color=DARK_BLUE, bold=True, font_name='Calibri')
edge_highlights = [
    "Runs on NVIDIA Jetson Nano / Raspberry Pi 4",
    "Model size: 25.4 MB (PT) / 12.7 MB (ONNX)",
    "No internet or cloud required for inference",
    "TensorRT & ONNX optimized for edge devices",
    "Factory floor ready — real-time 30 FPS processing",
]
add_multi_text(slide, Inches(7.1), Inches(5.05), Inches(5.5), Inches(1.8), edge_highlights,
               font_size=12, font_color=DARK_GRAY, bullet=True, line_spacing=1.5)

add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.35), SLIDE_WIDTH, Inches(0.15), fill_color=DARK_BLUE)
add_slide_number(slide, 7)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8: SYSTEM ARCHITECTURE
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_top_bar(slide)
add_section_header(slide, "SYSTEM ARCHITECTURE", "End-to-End Pipeline Design")

# Architecture flow diagram using boxes and arrows
arch_blocks = [
    ("Camera\nInput", Inches(0.3), Inches(2.2), Inches(1.5), Inches(1.0), MEDIUM_BLUE),
    ("Image\nCapture", Inches(2.2), Inches(2.2), Inches(1.5), Inches(1.0), NAVY_BLUE),
    ("Pre-\nprocessing", Inches(4.1), Inches(2.2), Inches(1.5), Inches(1.0), NAVY_BLUE),
    ("YOLOv8\nInference", Inches(6.0), Inches(2.2), Inches(1.8), Inches(1.0), MEDIUM_BLUE),
    ("Post-\nprocessing", Inches(8.2), Inches(2.2), Inches(1.5), Inches(1.0), NAVY_BLUE),
    ("Defect\nClassification", Inches(10.1), Inches(2.2), Inches(1.8), Inches(1.0), MEDIUM_BLUE),
]

for label, left, top, w, h, color in arch_blocks:
    add_card(slide, left, top, w, h, fill_color=color, border_color=LIGHT_BLUE)
    add_textbox(slide, left, top, w, h,
                label, font_size=11, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER, font_name='Calibri')

# Arrows
for i in range(5):
    positions = [Inches(1.8), Inches(3.7), Inches(5.6), Inches(7.8), Inches(9.7)]
    add_textbox(slide, positions[i], Inches(2.4), Inches(0.5), Inches(0.5),
                "→", font_size=20, font_color=LIGHT_BLUE, alignment=PP_ALIGN.CENTER)

# Output branches
outputs = [
    ("Alert System", Inches(8.5), Inches(3.7)),
    ("Web Dashboard", Inches(10.5), Inches(3.7)),
    ("Database", Inches(8.5), Inches(4.5)),
    ("Report Gen.", Inches(10.5), Inches(4.5)),
]

for label, x, y in outputs:
    add_card(slide, x, y, Inches(1.8), Inches(0.6), fill_color=NAVY_BLUE, border_color=ACCENT_BLUE)
    add_textbox(slide, x, y, Inches(1.8), Inches(0.6),
                label, font_size=10, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER, font_name='Calibri')

# Model architecture details
add_card(slide, Inches(0.3), Inches(3.7), Inches(7.8), Inches(3.2), fill_color=NAVY_BLUE, border_color=MEDIUM_BLUE)
add_textbox(slide, Inches(0.6), Inches(3.85), Inches(7.2), Inches(0.35),
            "MODEL ARCHITECTURE: YOLOv8_ResBlock_CBAM_GhostConv", font_size=14, font_color=LIGHT_BLUE, bold=True, font_name='Calibri')

model_arch = [
    "Backbone (GhostConv): Conv 3→16→32→64→128→256  |  2.1M params (51%)",
    "Neck (Feature Pyramid): SPPF + Upsampling + ResBlock_CBAM  |  1.5M params (37%)",
    "Detection Head: 6-class detector + BBox regression  |  0.46M params (11%)",
    "Total: 313 Layers  |  4.05M Parameters  |  10.2 GFLOPs  |  25.4 MB",
]
add_multi_text(slide, Inches(0.6), Inches(4.3), Inches(7.2), Inches(2.2), model_arch,
               font_size=12, font_color=OFF_WHITE, bullet=True, line_spacing=1.7)

add_bottom_bar(slide)
add_slide_number(slide, 8)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9: TECHNICAL IMPLEMENTATION
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_white_bg(slide)
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), fill_color=DARK_BLUE)
add_white_section_header(slide, "TECHNICAL IMPLEMENTATION", "Development Process & Methodology")

# Left column: Development Process
add_card(slide, Inches(0.5), Inches(1.7), Inches(6), Inches(5.4), fill_color=OFF_WHITE)
add_textbox(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(0.35),
            "DEVELOPMENT PROCESS", font_size=16, font_color=DARK_BLUE, bold=True, font_name='Calibri')

dev_steps = [
    "1. Data Preparation — NEU-DET dataset (1,800 images, 6 classes)",
    "2. Model Customization — YOLOv8s with ResCBAM + GhostConv",
    "3. Transfer Learning — COCO pre-trained weights fine-tuned",
    "4. Training — 100 epochs, batch size 16, AdamW optimizer",
    "5. Validation — mAP, precision, recall evaluation metrics",
    "6. Model Export — ONNX (12.7 MB) for edge deployment",
    "7. Web Integration — Flask dashboard with REST APIs",
    "8. Edge Deployment — Jetson Nano / Raspberry Pi ready",
]
add_multi_text(slide, Inches(0.8), Inches(2.3), Inches(5.5), Inches(4.5), dev_steps,
               font_size=12, font_color=DARK_GRAY, bullet=False, line_spacing=1.6)

# Right column: Training Configuration
add_card(slide, Inches(6.8), Inches(1.7), Inches(6), Inches(2.5), fill_color=DARK_BLUE)
add_textbox(slide, Inches(7.1), Inches(1.8), Inches(5.5), Inches(0.35),
            "TRAINING CONFIGURATION", font_size=16, font_color=LIGHT_BLUE, bold=True, font_name='Calibri')

config = [
    "Epochs: 100  |  Batch Size: 16  |  Image Size: 640×640",
    "Optimizer: AdamW (lr=0.001, momentum=0.937)",
    "Warmup: 3 epochs  |  Early Stopping: 20 patience",
    "GPU: NVIDIA RTX 3050 (4GB VRAM, CUDA 11.8)",
]
add_multi_text(slide, Inches(7.1), Inches(2.25), Inches(5.5), Inches(1.7), config,
               font_size=12, font_color=OFF_WHITE, bullet=True, line_spacing=1.6)

# AI/ML details
add_card(slide, Inches(6.8), Inches(4.4), Inches(6), Inches(2.7), fill_color=DARK_BLUE)
add_textbox(slide, Inches(7.1), Inches(4.5), Inches(5.5), Inches(0.35),
            "AI/ML APPROACH", font_size=16, font_color=LIGHT_BLUE, bold=True, font_name='Calibri')

ml_approach = [
    "Object Detection: YOLOv8 (You Only Look Once v8)",
    "Attention: CBAM (Channel + Spatial Attention)",
    "Augmentation: Mosaic, MixUp, HSV, Flip",
    "Loss: WIOU + Classification + Objectness",
    "Edge Optimization: Quantization, ONNX, TensorRT",
]
add_multi_text(slide, Inches(7.1), Inches(4.95), Inches(5.5), Inches(1.8), ml_approach,
               font_size=12, font_color=OFF_WHITE, bullet=True, line_spacing=1.5)

add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.35), SLIDE_WIDTH, Inches(0.15), fill_color=DARK_BLUE)
add_slide_number(slide, 9)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10: TECHNOLOGIES AND FRAMEWORKS
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_top_bar(slide)
add_section_header(slide, "TECHNOLOGIES & FRAMEWORKS", "Our Tech Stack")

# Tech categories in grid
tech_items = [
    ("💻", "PROGRAMMING", ["Python 3.13", "HTML5 / CSS3", "JavaScript"]),
    ("🧠", "AI FRAMEWORKS", ["Ultralytics YOLOv8", "PyTorch 2.7.1", "ONNX Runtime"]),
    ("📷", "COMPUTER VISION", ["OpenCV 4.8.0", "PIL / Pillow", "NumPy / SciPy"]),
    ("🌐", "WEB & API", ["Flask 2.3.0", "REST APIs", "Gunicorn"]),
    ("🔧", "EDGE HARDWARE", ["NVIDIA Jetson Nano", "Raspberry Pi 4", "NVIDIA RTX 3050"]),
    ("☁️", "CLOUD & DEPLOY", ["AWS IoT Core", "Docker", "TensorRT"]),
    ("💾", "DATABASES", ["SQLite (Edge)", "PostgreSQL (Cloud)", "YAML Config"]),
    ("🛠️", "DEV TOOLS", ["VS Code", "Git / GitHub", "Jupyter Notebook"]),
]

for i, (icon, title, techs) in enumerate(tech_items):
    row = i // 4
    col = i % 4
    x = Inches(0.5 + col * 3.15)
    y = Inches(1.7 + row * 2.7)
    
    add_card(slide, x, y, Inches(2.85), Inches(2.4), fill_color=NAVY_BLUE, border_color=MEDIUM_BLUE)
    add_textbox(slide, x, Inches(y.inches + 0.1), Inches(2.85), Inches(0.4),
                icon, font_size=24, font_color=WHITE, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(x.inches + 0.15), Inches(y.inches + 0.5), Inches(2.55), Inches(0.3),
                title, font_size=11, font_color=LIGHT_BLUE, bold=True, alignment=PP_ALIGN.CENTER, font_name='Calibri')
    add_multi_text(slide, Inches(x.inches + 0.3), Inches(y.inches + 0.85), Inches(2.25), Inches(1.3), techs,
                   font_size=11, font_color=OFF_WHITE, bullet=True, line_spacing=1.5)

add_bottom_bar(slide)
add_slide_number(slide, 10)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 11: CURRENT PROGRESS
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_white_bg(slide)
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), fill_color=DARK_BLUE)
add_white_section_header(slide, "CURRENT PROGRESS", "Development Milestones Achieved")

# Progress timeline
milestones = [
    ("✅", "Dataset Ready", "NEU-DET dataset\nprepared with 1,800\nimages across 6 classes"),
    ("✅", "Model Trained", "YOLOv8 custom model\ntrained for 100 epochs\nwith 79.2% mAP@50"),
    ("✅", "Web App Built", "Flask dashboard with\nupload, detection,\nand visualization"),
    ("✅", "Model Exported", "ONNX export (12.7MB)\nready for edge\ndeployment"),
    ("🔄", "Edge Deploy", "Hardware integration\nwith Jetson Nano\nin progress"),
]

for i, (status, title, desc) in enumerate(milestones):
    x = Inches(0.3 + i * 2.55)
    # Status icon
    add_textbox(slide, x, Inches(1.8), Inches(2.3), Inches(0.45),
                status, font_size=28, font_color=DARK_BLUE, alignment=PP_ALIGN.CENTER)
    # Card
    add_card(slide, x, Inches(2.35), Inches(2.3), Inches(2.3), fill_color=DARK_BLUE)
    add_textbox(slide, Inches(x.inches + 0.1), Inches(2.5), Inches(2.1), Inches(0.35),
                title, font_size=14, font_color=LIGHT_BLUE, bold=True, alignment=PP_ALIGN.CENTER, font_name='Calibri')
    add_textbox(slide, Inches(x.inches + 0.1), Inches(2.9), Inches(2.1), Inches(1.5),
                desc, font_size=11, font_color=OFF_WHITE, alignment=PP_ALIGN.CENTER, font_name='Calibri')

# Prototype status
add_card(slide, Inches(0.5), Inches(5.0), Inches(12.3), Inches(2.1), fill_color=OFF_WHITE)
add_textbox(slide, Inches(0.8), Inches(5.1), Inches(5), Inches(0.35),
            "PROTOTYPE STATUS", font_size=16, font_color=DARK_BLUE, bold=True, font_name='Calibri')

proto_items = [
    "Software prototype fully functional — Flask web app with real-time defect detection",
    "Custom YOLOv8 model achieving 79.2% mAP@50 — best among all compared models",
    "REST API endpoints operational: /upload, /api/model-info, /api/stats, /health",
    "Hardware integration with NVIDIA Jetson Nano — in progress (ONNX export ready)",
]
add_multi_text(slide, Inches(0.8), Inches(5.5), Inches(11.5), Inches(1.4), proto_items,
               font_size=12, font_color=DARK_GRAY, bullet=True, line_spacing=1.5)

add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.35), SLIDE_WIDTH, Inches(0.15), fill_color=DARK_BLUE)
add_slide_number(slide, 11)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 12: CHALLENGES FACED
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_top_bar(slide)
add_section_header(slide, "CHALLENGES FACED", "Obstacles Encountered & Solutions Applied")

challenges = [
    ("🗃️", "LIMITED DATASET", 
     "Challenge: Only 1,800 images available for training across 6 defect classes.",
     "Solution: Aggressive data augmentation — mosaic, mixup, HSV, flip, scale."),
    ("🔍", "SMALL DEFECT FEATURES",
     "Challenge: Original 200×200 images had very small defect features to detect.",
     "Solution: Upscaled to 640×640 pixels and used CBAM attention mechanism."),
    ("⚙️", "EDGE DEPLOYMENT",
     "Challenge: Full YOLOv8 too heavy for resource-constrained edge devices.",
     "Solution: GhostConv reduced parameters by 40%; ONNX export at 12.7MB."),
    ("⚖️", "CLASS IMBALANCE",
     "Challenge: Some defect types harder to detect, causing training imbalance.",
     "Solution: WIOU loss function with adjusted weights per class."),
    ("⏱️", "REAL-TIME SPEED",
     "Challenge: Need <100ms inference while maintaining detection accuracy.",
     "Solution: YOLOv8s (small variant) + TensorRT optimization achieves 35ms avg."),
    ("🔗", "HARDWARE INTEGRATION",
     "Challenge: Integrating AI model with industrial camera systems and IoT.",
     "Solution: Modular Flask API architecture with configurable input sources."),
]

for i, (icon, title, challenge, solution) in enumerate(challenges):
    row = i // 3
    col = i % 3
    x = Inches(0.3 + col * 4.25)
    y = Inches(1.7 + row * 2.7)
    
    add_card(slide, x, y, Inches(3.95), Inches(2.4), fill_color=NAVY_BLUE, border_color=MEDIUM_BLUE)
    add_textbox(slide, x, Inches(y.inches + 0.1), Inches(0.5), Inches(0.35),
                icon, font_size=20, font_color=WHITE)
    add_textbox(slide, Inches(x.inches + 0.5), Inches(y.inches + 0.1), Inches(3.25), Inches(0.3),
                title, font_size=12, font_color=LIGHT_BLUE, bold=True, font_name='Calibri')
    add_textbox(slide, Inches(x.inches + 0.15), Inches(y.inches + 0.55), Inches(3.65), Inches(0.7),
                challenge, font_size=10, font_color=OFF_WHITE, font_name='Calibri')
    add_textbox(slide, Inches(x.inches + 0.15), Inches(y.inches + 1.3), Inches(3.65), Inches(0.7),
                solution, font_size=10, font_color=GREEN, font_name='Calibri')

add_bottom_bar(slide)
add_slide_number(slide, 12)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 13: RESULTS AND ACHIEVEMENTS
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_white_bg(slide)
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), fill_color=DARK_BLUE)
add_white_section_header(slide, "RESULTS & ACHIEVEMENTS", "Performance Metrics & Benchmarks")

# Model comparison table
add_card(slide, Inches(0.5), Inches(1.7), Inches(7.5), Inches(3.5), fill_color=DARK_BLUE)
add_textbox(slide, Inches(0.8), Inches(1.8), Inches(7), Inches(0.35),
            "MODEL COMPARISON ON NEU-DET DATASET", font_size=14, font_color=LIGHT_BLUE, bold=True, font_name='Calibri')

# Table headers
headers = ["Model", "Params", "FLOPs", "mAP50", "mAP50-95"]
for j, h in enumerate(headers):
    add_textbox(slide, Inches(0.8 + j * 1.4), Inches(2.25), Inches(1.3), Inches(0.3),
                h, font_size=11, font_color=MEDIUM_BLUE, bold=True, alignment=PP_ALIGN.CENTER, font_name='Calibri')

# Table rows
rows = [
    ("YOLOv5s", "7.04M", "15.9G", "70.87%", "35.02%"),
    ("YOLOv5m", "20.89M", "183.5G", "71.74%", "36.67%"),
    ("YOLOv7", "37.22M", "196.2G", "71.92%", "37.2%"),
    ("YOLOv8", "3.01M", "8.1G", "77.7%", "46.2%"),
    ("Ours ★", "4.05M", "10.2G", "79.2%", "47.0%"),
]

for i, row_data in enumerate(rows):
    y = Inches(2.6 + i * 0.4)
    is_ours = i == 4
    fc = LIGHT_BLUE if is_ours else OFF_WHITE
    for j, val in enumerate(row_data):
        add_textbox(slide, Inches(0.8 + j * 1.4), y, Inches(1.3), Inches(0.3),
                    val, font_size=11, font_color=fc, bold=is_ours, alignment=PP_ALIGN.CENTER, font_name='Calibri')

# Performance improvements
add_card(slide, Inches(8.3), Inches(1.7), Inches(4.5), Inches(3.5), fill_color=DARK_BLUE)
add_textbox(slide, Inches(8.6), Inches(1.8), Inches(4), Inches(0.35),
            "PERFORMANCE vs MANUAL", font_size=14, font_color=LIGHT_BLUE, bold=True, font_name='Calibri')

perf_items = [
    ("Cost/Inspection", "$100-300 → $2-6", "98% ↓"),
    ("Time/Aircraft", "4 hours → 15 min", "94% ↓"),
    ("Accuracy", "85% → 99%", "+14%"),
    ("Throughput", "6/day → 96/day", "1500% ↑"),
    ("Annual Cost", "$200K → $50K", "75% ↓"),
]

for i, (metric, change, improvement) in enumerate(perf_items):
    y = Inches(2.3 + i * 0.5)
    add_textbox(slide, Inches(8.5), y, Inches(1.5), Inches(0.3),
                metric, font_size=10, font_color=OFF_WHITE, font_name='Calibri')
    add_textbox(slide, Inches(10.0), y, Inches(1.8), Inches(0.3),
                change, font_size=10, font_color=OFF_WHITE, font_name='Calibri')
    add_textbox(slide, Inches(11.8), y, Inches(0.8), Inches(0.3),
                improvement, font_size=10, font_color=GREEN, bold=True, font_name='Calibri')

# Ablation study
add_card(slide, Inches(0.5), Inches(5.5), Inches(12.3), Inches(1.6), fill_color=OFF_WHITE)
add_textbox(slide, Inches(0.8), Inches(5.6), Inches(5), Inches(0.35),
            "ABLATION STUDY RESULTS", font_size=14, font_color=DARK_BLUE, bold=True, font_name='Calibri')

ablation = [
    "YOLOv8 (base): 77.7%  →  +GhostConv: 78.7%  →  +ResCBAM: 78.2%  →  Ours (Combined): 79.2% mAP@50",
    "Precision improvement: +5.6%    |    F1-Score improvement: +3.1%    |    Inference Speed: 35ms average",
]
add_multi_text(slide, Inches(0.8), Inches(6.0), Inches(11.5), Inches(0.9), ablation,
               font_size=12, font_color=DARK_GRAY, bullet=True, line_spacing=1.5)

add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.35), SLIDE_WIDTH, Inches(0.15), fill_color=DARK_BLUE)
add_slide_number(slide, 13)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 14: DEMONSTRATION
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_top_bar(slide)
add_section_header(slide, "DEMONSTRATION", "Live System Prototype & Screenshots")

# Screenshot placeholders
add_card(slide, Inches(0.5), Inches(1.7), Inches(5.8), Inches(3.5), fill_color=NAVY_BLUE, border_color=MEDIUM_BLUE)
add_textbox(slide, Inches(0.8), Inches(1.8), Inches(5.3), Inches(0.35),
            "WEB DASHBOARD — UPLOAD INTERFACE", font_size=13, font_color=LIGHT_BLUE, bold=True, 
            alignment=PP_ALIGN.CENTER, font_name='Calibri')

# Try to add actual image
web_screenshot_path = os.path.join("imgs", "val_pred.jpg")
if os.path.exists(web_screenshot_path):
    slide.shapes.add_picture(web_screenshot_path, Inches(0.8), Inches(2.2), Inches(5.2), Inches(2.8))
else:
    add_textbox(slide, Inches(0.8), Inches(3.0), Inches(5.2), Inches(1.0),
                "[ Web Dashboard Screenshot ]\nFlask-based upload and detection interface\nwith drag-and-drop file support",
                font_size=14, font_color=OFF_WHITE, alignment=PP_ALIGN.CENTER, font_name='Calibri')

add_card(slide, Inches(6.6), Inches(1.7), Inches(6.2), Inches(3.5), fill_color=NAVY_BLUE, border_color=MEDIUM_BLUE)
add_textbox(slide, Inches(6.9), Inches(1.8), Inches(5.7), Inches(0.35),
            "DETECTION RESULTS — DEFECTS HIGHLIGHTED", font_size=13, font_color=LIGHT_BLUE, bold=True, 
            alignment=PP_ALIGN.CENTER, font_name='Calibri')

# Try to add confusion matrix or heatmap
heatmap_path = os.path.join("imgs", "heatmap.png")
if os.path.exists(heatmap_path):
    slide.shapes.add_picture(heatmap_path, Inches(6.9), Inches(2.2), Inches(5.6), Inches(2.8))
else:
    add_textbox(slide, Inches(6.9), Inches(3.0), Inches(5.6), Inches(1.0),
                "[ Detection Results Screenshot ]\nBounding boxes with confidence scores\nand defect classification labels",
                font_size=14, font_color=OFF_WHITE, alignment=PP_ALIGN.CENTER, font_name='Calibri')

# Features and links
add_card(slide, Inches(0.5), Inches(5.5), Inches(6), Inches(1.6), fill_color=NAVY_BLUE, border_color=MEDIUM_BLUE)
add_textbox(slide, Inches(0.8), Inches(5.6), Inches(5.5), Inches(0.3),
            "DEMO FEATURES", font_size=13, font_color=LIGHT_BLUE, bold=True, font_name='Calibri')
demo_features = [
    "Drag-and-drop image upload (JPG, PNG, BMP)",
    "Real-time defect detection with bounding boxes",
    "Color-coded classification with confidence scores",
]
add_multi_text(slide, Inches(0.8), Inches(5.95), Inches(5.5), Inches(1.0), demo_features,
               font_size=11, font_color=OFF_WHITE, bullet=True, line_spacing=1.4)

add_card(slide, Inches(6.8), Inches(5.5), Inches(6), Inches(1.6), fill_color=NAVY_BLUE, border_color=MEDIUM_BLUE)
add_textbox(slide, Inches(7.1), Inches(5.6), Inches(5.5), Inches(0.3),
            "PROJECT LINKS", font_size=13, font_color=LIGHT_BLUE, bold=True, font_name='Calibri')
links = [
    "GitHub: github.com/MUGESH-KUMAR-M/AeroVision-AI",
    "API Endpoints: /upload, /api/model-info, /health",
    "Local Demo: localhost:5000",
]
add_multi_text(slide, Inches(7.1), Inches(5.95), Inches(5.5), Inches(1.0), links,
               font_size=11, font_color=OFF_WHITE, bullet=True, line_spacing=1.4)

add_bottom_bar(slide)
add_slide_number(slide, 14)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 15: FUTURE ENHANCEMENTS
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_white_bg(slide)
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), fill_color=DARK_BLUE)
add_white_section_header(slide, "FUTURE ENHANCEMENTS", "Roadmap for Growth & Expansion")

# Enhancement categories
enhancements = [
    ("🚀", "ADDITIONAL FEATURES", [
        "Thermal/infrared imaging for subsurface defects",
        "AR overlay for technician guidance",
        "Multi-camera fusion for 360° inspection",
        "Predictive maintenance integration",
    ]),
    ("📈", "SCALABILITY PLANS", [
        "Multi-production line deployment",
        "Cloud SaaS platform launch",
        "Federated learning across facilities",
        "Multi-aircraft type support",
    ]),
    ("🏭", "COMMERCIAL DEPLOYMENT", [
        "Industry certification (DO-254)",
        "Enterprise licensing model",
        "Partnership with MRO facilities",
        "Global deployment infrastructure",
    ]),
]

for i, (icon, title, items) in enumerate(enhancements):
    x = Inches(0.3 + i * 4.25)
    add_card(slide, x, Inches(1.7), Inches(3.95), Inches(4.0), fill_color=DARK_BLUE)
    add_textbox(slide, x, Inches(1.85), Inches(3.95), Inches(0.5),
                icon, font_size=32, font_color=WHITE, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(x.inches + 0.2), Inches(2.4), Inches(3.55), Inches(0.35),
                title, font_size=14, font_color=LIGHT_BLUE, bold=True, alignment=PP_ALIGN.CENTER, font_name='Calibri')
    add_multi_text(slide, Inches(x.inches + 0.3), Inches(2.9), Inches(3.35), Inches(2.5), items,
                   font_size=12, font_color=OFF_WHITE, bullet=True, line_spacing=1.6)

# Mobile & drone future
add_card(slide, Inches(0.5), Inches(6.0), Inches(12.3), Inches(1.1), fill_color=OFF_WHITE)
add_textbox(slide, Inches(0.8), Inches(6.15), Inches(11.8), Inches(0.7),
            "🔮  Future Vision:  Autonomous inspection drones with onboard AI  •  Mobile apps (Flutter/React Native)  •  "
            "IoT sensor fusion  •  Real-time video stream analysis  •  Enterprise suite with analytics dashboard",
            font_size=13, font_color=DARK_BLUE, alignment=PP_ALIGN.CENTER, font_name='Calibri')

add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.35), SLIDE_WIDTH, Inches(0.15), fill_color=DARK_BLUE)
add_slide_number(slide, 15)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 16: PROTOTYPE DEVELOPMENT PLAN
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_top_bar(slide)
add_section_header(slide, "PROTOTYPE DEVELOPMENT PLAN", "3-Month Development Timeline")

months = [
    ("MONTH 1", "Foundation & Data", [
        "Hardware procurement & setup",
        "Industrial camera integration",
        "NEU-DET dataset preparation",
        "Data augmentation pipeline",
        "Development environment config",
    ], MEDIUM_BLUE),
    ("MONTH 2", "Model & Testing", [
        "YOLOv8 model customization",
        "ResCBAM + GhostConv integration",
        "100-epoch training cycle",
        "Validation & hyperparameter tuning",
        "Flask web app development",
    ], ACCENT_BLUE),
    ("MONTH 3", "Prototype & Validation", [
        "Physical prototype assembly",
        "Edge device deployment (Jetson)",
        "End-to-end system testing",
        "Performance benchmarking",
        "Documentation & demo prep",
    ], SKY_BLUE),
]

for i, (month, title, tasks, color) in enumerate(months):
    x = Inches(0.3 + i * 4.25)
    
    # Month header
    add_card(slide, x, Inches(1.7), Inches(3.95), Inches(0.65), fill_color=color)
    add_textbox(slide, x, Inches(1.75), Inches(3.95), Inches(0.55),
                f"{month}: {title}", font_size=16, font_color=WHITE, bold=True, 
                alignment=PP_ALIGN.CENTER, font_name='Calibri')
    
    # Tasks card
    add_card(slide, x, Inches(2.4), Inches(3.95), Inches(3.2), fill_color=NAVY_BLUE, border_color=color)
    add_multi_text(slide, Inches(x.inches + 0.2), Inches(2.6), Inches(3.55), Inches(2.8), tasks,
                   font_size=12, font_color=OFF_WHITE, bullet=True, line_spacing=1.6)

# Progress indicator
add_card(slide, Inches(0.5), Inches(6.0), Inches(12.3), Inches(1.1), fill_color=NAVY_BLUE, border_color=MEDIUM_BLUE)
add_textbox(slide, Inches(0.8), Inches(6.1), Inches(11.5), Inches(0.4),
            "CURRENT STATUS: Month 2 Complete — Model trained (79.2% mAP@50), Web app functional, ONNX export ready",
            font_size=14, font_color=GREEN, bold=True, alignment=PP_ALIGN.CENTER, font_name='Calibri')
add_textbox(slide, Inches(0.8), Inches(6.5), Inches(11.5), Inches(0.4),
            "Next: Physical prototype assembly & edge device deployment with NVIDIA Jetson Nano",
            font_size=12, font_color=OFF_WHITE, alignment=PP_ALIGN.CENTER, font_name='Calibri')

add_bottom_bar(slide)
add_slide_number(slide, 16)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 17: IMPACT AND SCALABILITY
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_white_bg(slide)
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), fill_color=DARK_BLUE)
add_white_section_header(slide, "IMPACT & SCALABILITY", "Creating Meaningful Change Across Industries")

# Impact categories
impacts = [
    ("🛡️", "SOCIAL IMPACT", [
        "Enhanced aviation safety for passengers",
        "Reduced human exposure to hazardous inspection environments",
        "Job upskilling: inspectors become AI operators",
        "Improved public confidence in air travel safety",
    ], MEDIUM_BLUE),
    ("🏭", "INDUSTRY IMPACT", [
        "70% reduction in inspection time",
        "60% cost reduction per facility",
        "1500% throughput increase (6→96/day)",
        "Standardized quality across all facilities",
    ], DARK_BLUE),
    ("🌿", "ENVIRONMENTAL BENEFITS", [
        "Reduced material waste from early detection",
        "Lower energy consumption vs cloud AI",
        "Extended aircraft lifecycle through maintenance",
        "Minimized hazardous chemical usage",
    ], GREEN),
    ("💰", "REVENUE OPPORTUNITIES", [
        "SaaS licensing: $50K-100K/year per facility",
        "Hardware + software bundle sales",
        "Per-inspection usage fees",
        "Consulting & customization services",
    ], ACCENT_BLUE),
]

for i, (icon, title, items, color) in enumerate(impacts):
    row = i // 2
    col = i % 2
    x = Inches(0.5 + col * 6.4)
    y = Inches(1.7 + row * 2.7)
    
    add_card(slide, x, y, Inches(6.1), Inches(2.4), fill_color=DARK_BLUE)
    add_textbox(slide, x, Inches(y.inches + 0.1), Inches(0.5), Inches(0.35),
                icon, font_size=22, font_color=WHITE)
    add_textbox(slide, Inches(x.inches + 0.55), Inches(y.inches + 0.1), Inches(5.2), Inches(0.3),
                title, font_size=14, font_color=color, bold=True, font_name='Calibri')
    add_multi_text(slide, Inches(x.inches + 0.55), Inches(y.inches + 0.5), Inches(5.2), Inches(1.7), items,
                   font_size=12, font_color=OFF_WHITE, bullet=True, line_spacing=1.5)

# Scalability note
add_card(slide, Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.85), fill_color=OFF_WHITE)
add_textbox(slide, Inches(0.8), Inches(6.4), Inches(11.8), Inches(0.6),
            "🌐  Scalability: Deploy across multiple production lines → MRO facilities → Airports → Defense → "
            "Industrial manufacturing → Automotive — same AI, different domains.",
            font_size=13, font_color=DARK_BLUE, alignment=PP_ALIGN.CENTER, font_name='Calibri')

add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.35), SLIDE_WIDTH, Inches(0.15), fill_color=DARK_BLUE)
add_slide_number(slide, 17)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 18: CONCLUSION
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_top_bar(slide)
add_section_header(slide, "CONCLUSION", "Summary of AeroVision AI")

# Summary
add_card(slide, Inches(0.5), Inches(1.7), Inches(12.3), Inches(1.5), fill_color=NAVY_BLUE, border_color=MEDIUM_BLUE)
add_textbox(slide, Inches(0.8), Inches(1.8), Inches(11.8), Inches(1.2),
            "AeroVision AI transforms aircraft surface inspection from a slow, error-prone manual process into a "
            "fast, accurate, and cost-effective AI-powered system. Using a custom YOLOv8 model with ResCBAM attention "
            "and GhostConv, deployed on edge devices, we achieve real-time multi-defect detection that sets a new "
            "benchmark for aerospace quality inspection.",
            font_size=16, font_color=WHITE, alignment=PP_ALIGN.CENTER, font_name='Calibri')

# Key benefits
benefits = [
    ("79.2%", "mAP@50\nAccuracy"),
    ("35ms", "Average\nInference"),
    ("4.05M", "Parameters\n(Lightweight)"),
    ("98%", "Cost\nReduction"),
    ("6", "Defect\nTypes"),
    ("$486K", "5-Year\nROI"),
]

for i, (value, label) in enumerate(benefits):
    x = Inches(0.3 + i * 2.15)
    add_card(slide, x, Inches(3.5), Inches(1.95), Inches(1.6), fill_color=NAVY_BLUE, border_color=MEDIUM_BLUE)
    add_textbox(slide, x, Inches(3.6), Inches(1.95), Inches(0.55),
                value, font_size=28, font_color=MEDIUM_BLUE, bold=True, alignment=PP_ALIGN.CENTER, font_name='Calibri')
    add_textbox(slide, x, Inches(4.15), Inches(1.95), Inches(0.7),
                label, font_size=11, font_color=OFF_WHITE, alignment=PP_ALIGN.CENTER, font_name='Calibri')

# Future outcomes
add_card(slide, Inches(0.5), Inches(5.4), Inches(12.3), Inches(1.7), fill_color=NAVY_BLUE, border_color=MEDIUM_BLUE)
add_textbox(slide, Inches(0.8), Inches(5.5), Inches(5), Inches(0.35),
            "EXPECTED FUTURE OUTCOMES", font_size=14, font_color=LIGHT_BLUE, bold=True, font_name='Calibri')

future = [
    "Industry certification and commercial deployment within 12 months",
    "Expansion to automotive & industrial heavy machinery inspection",
    "Autonomous drone-based inspection integration",
    "SaaS platform serving 100+ aerospace facilities globally",
]
add_multi_text(slide, Inches(0.8), Inches(5.9), Inches(11.5), Inches(1.0), future,
               font_size=12, font_color=OFF_WHITE, bullet=True, line_spacing=1.4)

add_bottom_bar(slide)
add_slide_number(slide, 18)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 19: THANK YOU
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)

# Left accent bar
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.15), SLIDE_HEIGHT, fill_color=MEDIUM_BLUE)

# Top bar
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.08), fill_color=MEDIUM_BLUE)

# Main thank you text
add_textbox(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.2),
            "THANK YOU", font_size=60, font_color=WHITE, bold=True, 
            alignment=PP_ALIGN.CENTER, font_name='Calibri')

add_textbox(slide, Inches(1), Inches(2.7), Inches(11), Inches(0.6),
            "AeroVision AI — Redefining Aerospace Inspection with AI at the Edge",
            font_size=20, font_color=LIGHT_BLUE, alignment=PP_ALIGN.CENTER, font_name='Calibri')

# Separator
add_horizontal_line(slide, Inches(4), Inches(3.5), Inches(5.3), color=MEDIUM_BLUE, thickness=3)

# Contact details
add_card(slide, Inches(2.5), Inches(4.0), Inches(8.3), Inches(2.8), fill_color=NAVY_BLUE, border_color=MEDIUM_BLUE)

contact_items = [
    ("👤", "Team Lead: Mugesh Kumar M"),
    ("👥", "Team: Karthikeyan S  |  Gokulakrishnan S  |  Pradeep Kumar R"),
    ("🎓", "K.S.Rangasamy College of Technology, Tiruchengode"),
    ("📧", "Email: mugeshkumarm.22ai@ksr.ac.in"),
    ("🔗", "GitHub: github.com/MUGESH-KUMAR-M/AeroVision-AI"),
]

for i, (icon, text) in enumerate(contact_items):
    y = Inches(4.2 + i * 0.45)
    add_textbox(slide, Inches(3.5), y, Inches(0.5), Inches(0.35),
                icon, font_size=16, font_color=WHITE)
    add_textbox(slide, Inches(4.1), y, Inches(6), Inches(0.35),
                text, font_size=14, font_color=OFF_WHITE, font_name='Calibri')

# Bottom competition badge
add_textbox(slide, Inches(3), Inches(6.7), Inches(7.3), Inches(0.5),
            "Tata Technologies InnoVent 2026  •  Aerospace & Aviation  •  AI-at-the-Edge",
            font_size=12, font_color=MEDIUM_BLUE, alignment=PP_ALIGN.CENTER, font_name='Calibri')

# Bottom bar
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.35), SLIDE_WIDTH, Inches(0.15), fill_color=MEDIUM_BLUE)
add_slide_number(slide, 19)


# ═══════════════════════════════════════════════════════════════════════════════
# SAVE THE PRESENTATION
# ═══════════════════════════════════════════════════════════════════════════════
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "AeroVision_AI_Tata_InnoVent.pptx")
prs.save(output_path)
print(f"\n✅ Presentation saved successfully!")
print(f"📁 File: {output_path}")
print(f"📊 Total slides: {len(prs.slides)}")
print(f"📐 Format: 16:9 Widescreen")
print(f"🎨 Theme: White & Dark Blue")
